import io
import json
import sqlite3
import streamlit as st
import plotly.graph_objects as go
import plotly.express as px
import anthropic
from pptx import Presentation
import pdfplumber

# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Strategy Evaluator",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── CSS ───────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
  .hero {
    background: linear-gradient(135deg, #1e3a5f 0%, #2d6a9f 100%);
    padding: 2rem 2.5rem; border-radius: 14px; color: white; margin-bottom: 1.8rem;
  }
  .hero h1 { margin: 0; font-size: 2rem; }
  .hero p  { margin: 0.4rem 0 0; opacity: .85; }
  .kpi-box {
    background: white; border: 1px solid #e0e7f0; border-radius: 10px;
    padding: 1.1rem; text-align: center; box-shadow: 0 2px 6px rgba(0,0,0,.05);
  }
  .kpi-box .val { font-size: 2rem; font-weight: 700; }
  .kpi-box .lbl { font-size: .8rem; color: #666; margin-top: .2rem; }
  .brand-row {
    background: #f4f8fd; border-left: 4px solid #2d6a9f;
    padding: .9rem 1rem; border-radius: 0 8px 8px 0; margin-bottom: .8rem;
  }
  .verdict-chip {
    display: inline-block; padding: .25rem .9rem; border-radius: 20px;
    font-weight: 700; font-size: .8rem;
  }
  .rec-item {
    padding: .55rem .8rem .55rem 1rem; margin: .45rem 0;
    border-radius: 0 6px 6px 0;
  }
</style>
""", unsafe_allow_html=True)

# ══════════════════════════════════════════════════════════════════════════════
# DATABASE
# ══════════════════════════════════════════════════════════════════════════════
DB_PATH = "strategy_data.db"


def get_conn():
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    conn = get_conn()
    conn.execute("""
        CREATE TABLE IF NOT EXISTS brands (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL UNIQUE,
            industry TEXT NOT NULL,
            description TEXT,
            target_market TEXT,
            created_at TEXT DEFAULT CURRENT_TIMESTAMP
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS analyses (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            brand_id INTEGER NOT NULL,
            deck_name TEXT NOT NULL,
            overall_score INTEGER,
            verdict TEXT,
            results_json TEXT,
            created_at TEXT DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (brand_id) REFERENCES brands(id)
        )
    """)
    conn.commit()
    conn.close()


def add_brand(name, industry, description, target_market):
    conn = get_conn()
    try:
        conn.execute(
            "INSERT INTO brands (name, industry, description, target_market) VALUES (?, ?, ?, ?)",
            (name, industry, description, target_market),
        )
        conn.commit()
        return True
    except sqlite3.IntegrityError:
        return False
    finally:
        conn.close()


def get_brands():
    conn = get_conn()
    rows = conn.execute("SELECT * FROM brands ORDER BY name").fetchall()
    conn.close()
    return [dict(r) for r in rows]


def get_brand(brand_id):
    conn = get_conn()
    row = conn.execute("SELECT * FROM brands WHERE id = ?", (brand_id,)).fetchone()
    conn.close()
    return dict(row) if row else None


def update_brand(brand_id, name, industry, description, target_market):
    conn = get_conn()
    try:
        conn.execute(
            "UPDATE brands SET name=?, industry=?, description=?, target_market=? WHERE id=?",
            (name, industry, description, target_market, brand_id),
        )
        conn.commit()
        return True
    except sqlite3.IntegrityError:
        return False
    finally:
        conn.close()


def delete_brand(brand_id):
    conn = get_conn()
    conn.execute("DELETE FROM analyses WHERE brand_id = ?", (brand_id,))
    conn.execute("DELETE FROM brands WHERE id = ?", (brand_id,))
    conn.commit()
    conn.close()


def save_analysis(brand_id, deck_name, overall_score, verdict, results_json):
    conn = get_conn()
    cursor = conn.execute(
        "INSERT INTO analyses (brand_id, deck_name, overall_score, verdict, results_json) VALUES (?, ?, ?, ?, ?)",
        (brand_id, deck_name, overall_score, verdict, json.dumps(results_json)),
    )
    aid = cursor.lastrowid
    conn.commit()
    conn.close()
    return aid


def get_analyses(brand_id=None):
    conn = get_conn()
    if brand_id:
        rows = conn.execute(
            """SELECT a.*, b.name as brand_name, b.industry
               FROM analyses a JOIN brands b ON a.brand_id = b.id
               WHERE a.brand_id = ? ORDER BY a.created_at DESC""",
            (brand_id,),
        ).fetchall()
    else:
        rows = conn.execute(
            """SELECT a.*, b.name as brand_name, b.industry
               FROM analyses a JOIN brands b ON a.brand_id = b.id
               ORDER BY a.created_at DESC"""
        ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


def delete_analysis(analysis_id):
    conn = get_conn()
    conn.execute("DELETE FROM analyses WHERE id = ?", (analysis_id,))
    conn.commit()
    conn.close()


# ══════════════════════════════════════════════════════════════════════════════
# FILE PARSING
# ══════════════════════════════════════════════════════════════════════════════
def extract_text_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def extract_text_from_pdf(file_bytes):
    parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                parts.append(f"[Page {i}]\n{text.strip()}")
    return "\n\n".join(parts)


def extract_text(file_bytes, filename):
    name = filename.lower()
    if name.endswith(".pptx"):
        return extract_text_from_pptx(file_bytes)
    elif name.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    raise ValueError(f"Unsupported file type: {filename}")


# ══════════════════════════════════════════════════════════════════════════════
# AI ANALYSIS
# ══════════════════════════════════════════════════════════════════════════════
def analyze_strategy(brand_name, industry, brand_description, target_market, deck_content, deck_name):
    client = anthropic.Anthropic(api_key=st.secrets["ANTHROPIC_API_KEY"])

    prompt = f"""You are a senior business strategy consultant. Analyze this strategy document for "{brand_name}" and evaluate its effectiveness.

Brand: {brand_name} | Industry: {industry}
Description: {brand_description or "Not provided"}
Target Market: {target_market or "Not provided"}
Document: {deck_name}

Strategy Content:
{deck_content[:15000]}

Score across 8 dimensions (0-100 each). Be specific and honest. Return ONLY valid JSON, no markdown:
{{
  "overall_score": <integer 0-100>,
  "executive_summary": "<2-3 sentences>",
  "verdict": "<Highly Effective | Effective | Needs Improvement | Ineffective>",
  "dimensions": {{
    "clarity": {{"score": <0-100>, "label": "Clarity", "feedback": "<2-3 sentences>", "strengths": ["<s>"], "improvements": ["<i>"]}},
    "target_audience": {{"score": <0-100>, "label": "Target Audience", "feedback": "<2-3 sentences>", "strengths": ["<s>"], "improvements": ["<i>"]}},
    "competitive_advantage": {{"score": <0-100>, "label": "Competitive Advantage", "feedback": "<2-3 sentences>", "strengths": ["<s>"], "improvements": ["<i>"]}},
    "feasibility": {{"score": <0-100>, "label": "Feasibility", "feedback": "<2-3 sentences>", "strengths": ["<s>"], "improvements": ["<i>"]}},
    "risk_assessment": {{"score": <0-100>, "label": "Risk Assessment", "feedback": "<2-3 sentences>", "strengths": ["<s>"], "improvements": ["<i>"]}},
    "kpis_metrics": {{"score": <0-100>, "label": "KPIs & Metrics", "feedback": "<2-3 sentences>", "strengths": ["<s>"], "improvements": ["<i>"]}},
    "market_alignment": {{"score": <0-100>, "label": "Market Alignment", "feedback": "<2-3 sentences>", "strengths": ["<s>"], "improvements": ["<i>"]}},
    "brand_consistency": {{"score": <0-100>, "label": "Brand Consistency", "feedback": "<2-3 sentences>", "strengths": ["<s>"], "improvements": ["<i>"]}}}},
  "top_strengths": ["<s1>", "<s2>", "<s3>"],
  "critical_risks": ["<r1>", "<r2>", "<r3>"],
  "recommendations": [
    {{"priority": "High", "action": "<step>"}},
    {{"priority": "High", "action": "<step>"}},
    {{"priority": "Medium", "action": "<step>"}},
    {{"priority": "Medium", "action": "<step>"}},
    {{"priority": "Low", "action": "<step>"}}
  ]
}}"""

    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=4096,
        messages=[{"role": "user", "content": prompt}],
    )
    text = msg.content[0].text.strip()
    if text.startswith("```"):
        lines = text.split("\n")
        text = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])
    return json.loads(text)


# ══════════════════════════════════════════════════════════════════════════════
# CHART HELPERS
# ══════════════════════════════════════════════════════════════════════════════
def score_color(s):
    if s >= 80: return "#27ae60"
    if s >= 65: return "#2980b9"
    if s >= 50: return "#f39c12"
    return "#e74c3c"


def verdict_color(v):
    return {"Highly Effective": "#27ae60", "Effective": "#2980b9",
            "Needs Improvement": "#f39c12", "Ineffective": "#e74c3c"}.get(v, "#666")


def radar_chart(dims):
    labels = [d["label"] for d in dims] + [dims[0]["label"]]
    scores = [d["score"] for d in dims] + [dims[0]["score"]]
    fig = go.Figure(go.Scatterpolar(
        r=scores, theta=labels, fill="toself",
        fillcolor="rgba(45,106,159,.18)",
        line=dict(color="#2d6a9f", width=2),
    ))
    fig.update_layout(
        polar=dict(radialaxis=dict(visible=True, range=[0, 100])),
        showlegend=False, height=390,
        margin=dict(l=55, r=55, t=35, b=35),
    )
    return fig


def bar_chart(dims):
    fig = go.Figure(go.Bar(
        x=[d["score"] for d in dims],
        y=[d["label"] for d in dims],
        orientation="h",
        marker_color=[score_color(d["score"]) for d in dims],
        text=[f"{d['score']}/100" for d in dims],
        textposition="outside",
    ))
    fig.update_layout(
        xaxis=dict(range=[0, 115], showgrid=True, gridcolor="#f0f0f0"),
        yaxis=dict(autorange="reversed"),
        height=340, plot_bgcolor="white", paper_bgcolor="white",
        margin=dict(l=10, r=55, t=15, b=15),
    )
    return fig


def gauge_chart(score):
    fig = go.Figure(go.Indicator(
        mode="gauge+number", value=score,
        gauge={
            "axis": {"range": [0, 100]},
            "bar": {"color": score_color(score)},
            "steps": [
                {"range": [0, 50],   "color": "#fde8e8"},
                {"range": [50, 65],  "color": "#fef3cd"},
                {"range": [65, 80],  "color": "#d1ecf1"},
                {"range": [80, 100], "color": "#d4edda"},
            ],
        },
    ))
    fig.update_layout(height=230, margin=dict(l=20, r=20, t=30, b=10))
    return fig


def render_results(results):
    vc = verdict_color(results["verdict"])
    c1, c2 = st.columns([1, 2])
    with c1:
        st.plotly_chart(gauge_chart(results["overall_score"]), use_container_width=True)
        st.markdown(
            f"<div style='text-align:center'><span class='verdict-chip' "
            f"style='background:{vc}18;color:{vc};border:1px solid {vc}'>"
            f"{results['verdict']}</span></div>",
            unsafe_allow_html=True,
        )
    with c2:
        st.markdown("**Executive Summary**")
        st.markdown(results["executive_summary"])
        ca, cb = st.columns(2)
        with ca:
            st.markdown("**Top Strengths**")
            for s in results.get("top_strengths", []):
                st.markdown(f"✅ {s}")
        with cb:
            st.markdown("**Critical Risks**")
            for r in results.get("critical_risks", []):
                st.markdown(f"⚠️ {r}")

    st.markdown("---")
    dims_list = [{"label": v["label"], "score": v["score"]} for v in results["dimensions"].values()]
    rc1, rc2 = st.columns(2)
    with rc1:
        st.subheader("Dimension Radar")
        st.plotly_chart(radar_chart(dims_list), use_container_width=True)
    with rc2:
        st.subheader("Dimension Scores")
        st.plotly_chart(bar_chart(dims_list), use_container_width=True)

    st.markdown("---")
    st.subheader("Detailed Breakdown")
    dim_items = list(results["dimensions"].items())
    for i in range(0, len(dim_items), 2):
        cols = st.columns(2)
        for j, col in enumerate(cols):
            if i + j >= len(dim_items):
                break
            _, dim = dim_items[i + j]
            sc = score_color(dim["score"])
            with col:
                with st.expander(f"{dim['label']} — {dim['score']}/100", expanded=True):
                    st.markdown(
                        f"**Score:** <span style='color:{sc};font-weight:700'>{dim['score']}/100</span>",
                        unsafe_allow_html=True,
                    )
                    st.markdown(dim["feedback"])
                    if dim.get("strengths"):
                        st.markdown("**Strengths:**")
                        for s in dim["strengths"]:
                            st.markdown(f"✅ {s}")
                    if dim.get("improvements"):
                        st.markdown("**Improvements:**")
                        for imp in dim["improvements"]:
                            st.markdown(f"🔧 {imp}")

    st.markdown("---")
    st.subheader("Recommendations")
    pc = {"High": "#e74c3c", "Medium": "#f39c12", "Low": "#27ae60"}
    for rec in results.get("recommendations", []):
        c = pc.get(rec["priority"], "#666")
        st.markdown(
            f"<div class='rec-item' style='border-left:3px solid {c};background:{c}10'>"
            f"<span style='color:{c};font-weight:700;font-size:.78rem'>{rec['priority'].upper()} PRIORITY</span><br>"
            f"{rec['action']}</div>",
            unsafe_allow_html=True,
        )


# ══════════════════════════════════════════════════════════════════════════════
# APP INIT
# ══════════════════════════════════════════════════════════════════════════════
init_db()

INDUSTRIES = [
    "Technology", "Retail & E-commerce", "Food & Beverage",
    "Fashion & Apparel", "Healthcare", "Finance & Banking",
    "Real Estate", "Education", "Entertainment & Media",
    "Automotive", "Travel & Hospitality", "Beauty & Personal Care",
    "Sports & Fitness", "Manufacturing", "Other",
]

# ── Sidebar ───────────────────────────────────────────────────────────────────
st.sidebar.markdown("## 📊 Strategy Evaluator")
st.sidebar.markdown("---")
page = st.sidebar.radio(
    "Navigate",
    ["🏠 Dashboard", "🏢 Brands", "🔍 Analyze", "📋 History"],
    label_visibility="collapsed",
)
st.sidebar.markdown("---")
_brands   = get_brands()
_analyses = get_analyses()
st.sidebar.markdown(f"**{len(_brands)}** brand(s) · **{len(_analyses)}** analysis(es)")


# ══════════════════════════════════════════════════════════════════════════════
# DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
if page == "🏠 Dashboard":
    st.markdown(
        "<div class='hero'><h1>📊 Strategy Evaluator</h1>"
        "<p>AI-powered strategy effectiveness analysis for your brands</p></div>",
        unsafe_allow_html=True,
    )

    avg_score = sum(a["overall_score"] for a in _analyses) / len(_analyses) if _analyses else 0
    effective = sum(1 for a in _analyses if a["verdict"] in ("Highly Effective", "Effective"))

    c1, c2, c3, c4 = st.columns(4)
    for col, val, lbl, color in [
        (c1, len(_brands),        "Brands",               "#2d6a9f"),
        (c2, len(_analyses),      "Analyses Run",          "#2d6a9f"),
        (c3, f"{avg_score:.0f}",  "Avg Score",             score_color(avg_score) if _analyses else "#999"),
        (c4, effective,           "Effective Strategies",  "#27ae60"),
    ]:
        col.markdown(
            f"<div class='kpi-box'><div class='val' style='color:{color}'>{val}</div>"
            f"<div class='lbl'>{lbl}</div></div>",
            unsafe_allow_html=True,
        )

    st.markdown("---")

    if not _analyses:
        st.info("No analyses yet. Add brands → upload a deck → run your first analysis.")
    else:
        cl, cr = st.columns([2, 1])
        with cl:
            st.subheader("Recent Analyses")
            for a in _analyses[:8]:
                sc = score_color(a["overall_score"])
                vc = verdict_color(a["verdict"])
                st.markdown(
                    f"<div class='brand-row'><strong>{a['deck_name']}</strong> · <em>{a['brand_name']}</em>"
                    f"<span style='float:right;color:{sc};font-weight:700'>{a['overall_score']}/100</span><br>"
                    f"<small style='color:#666'>{a['industry']} · {a['created_at'][:16]}</small>"
                    f"<span style='float:right;color:{vc};font-size:.8rem'>{a['verdict']}</span></div>",
                    unsafe_allow_html=True,
                )
        with cr:
            st.subheader("Score Distribution")
            buckets = {"0–49": 0, "50–64": 0, "65–79": 0, "80–100": 0}
            for a in _analyses:
                s = a["overall_score"]
                if s < 50:   buckets["0–49"] += 1
                elif s < 65: buckets["50–64"] += 1
                elif s < 80: buckets["65–79"] += 1
                else:        buckets["80–100"] += 1
            fig = px.pie(
                values=list(buckets.values()), names=list(buckets.keys()),
                color_discrete_sequence=["#e74c3c", "#f39c12", "#2980b9", "#27ae60"], hole=.35,
            )
            fig.update_layout(height=260, margin=dict(l=0, r=0, t=10, b=10))
            st.plotly_chart(fig, use_container_width=True)

    if _brands:
        st.markdown("---")
        st.subheader("Brand Overview")
        for brand in _brands:
            ba = [a for a in _analyses if a["brand_name"] == brand["name"]]
            avg = sum(a["overall_score"] for a in ba) / len(ba) if ba else None
            ca, cb, cc = st.columns([3, 1, 1])
            ca.markdown(f"**{brand['name']}** · {brand['industry']}")
            cb.markdown(f"Analyses: **{len(ba)}**")
            if avg is not None:
                cc.markdown(f"Avg: <span style='color:{score_color(avg)};font-weight:700'>{avg:.0f}</span>", unsafe_allow_html=True)
            else:
                cc.markdown("No analyses yet")


# ══════════════════════════════════════════════════════════════════════════════
# BRANDS
# ══════════════════════════════════════════════════════════════════════════════
elif page == "🏢 Brands":
    st.title("Brand Management")
    tab_add, tab_all = st.tabs(["➕ Add Brand", "📋 All Brands"])

    with tab_add:
        st.subheader("Register a New Brand")
        with st.form("add_brand"):
            c1, c2 = st.columns(2)
            with c1:
                b_name     = st.text_input("Brand Name *", placeholder="e.g., Acme Corp")
                b_industry = st.selectbox("Industry *", INDUSTRIES)
            with c2:
                b_market = st.text_input("Target Market", placeholder="e.g., Young professionals 25–35")
                b_desc   = st.text_area("Brand Description", placeholder="Values, positioning…", height=100)
            if st.form_submit_button("Add Brand", type="primary"):
                if not b_name.strip():
                    st.error("Brand name is required.")
                elif add_brand(b_name.strip(), b_industry, b_desc, b_market):
                    st.success(f"Brand '{b_name}' added!")
                    st.rerun()
                else:
                    st.error(f"'{b_name}' already exists.")

    with tab_all:
        brands = get_brands()
        if not brands:
            st.info("No brands yet.")
        for brand in brands:
            with st.expander(f"🏢 {brand['name']} — {brand['industry']}"):
                ekey = f"edit_{brand['id']}"
                if st.session_state.get(ekey):
                    with st.form(f"edit_{brand['id']}_form"):
                        ec1, ec2 = st.columns(2)
                        with ec1:
                            nn = st.text_input("Brand Name", value=brand["name"])
                            ni = st.selectbox("Industry", INDUSTRIES,
                                              index=INDUSTRIES.index(brand["industry"]) if brand["industry"] in INDUSTRIES else 0)
                        with ec2:
                            nm = st.text_input("Target Market", value=brand.get("target_market") or "")
                            nd = st.text_area("Description", value=brand.get("description") or "", height=80)
                        sc1, sc2 = st.columns(2)
                        if sc1.form_submit_button("Save", type="primary"):
                            update_brand(brand["id"], nn, ni, nd, nm)
                            st.session_state[ekey] = False
                            st.rerun()
                        if sc2.form_submit_button("Cancel"):
                            st.session_state[ekey] = False
                            st.rerun()
                else:
                    dc1, dc2 = st.columns([3, 1])
                    with dc1:
                        st.markdown(f"**Industry:** {brand['industry']}")
                        if brand.get("target_market"): st.markdown(f"**Target Market:** {brand['target_market']}")
                        if brand.get("description"):   st.markdown(f"**Description:** {brand['description']}")
                        st.markdown(f"**Added:** {brand['created_at'][:10]}")
                        st.markdown(f"**Analyses:** {len(get_analyses(brand['id']))}")
                    with dc2:
                        if st.button("✏️ Edit",   key=f"btn_edit_{brand['id']}"):
                            st.session_state[ekey] = True
                            st.rerun()
                        if st.button("🗑️ Delete", key=f"btn_del_{brand['id']}"):
                            delete_brand(brand["id"])
                            st.success("Brand deleted.")
                            st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
# ANALYZE
# ══════════════════════════════════════════════════════════════════════════════
elif page == "🔍 Analyze":
    st.title("Analyze Strategy Deck")
    brands = get_brands()
    if not brands:
        st.warning("Add at least one brand first.")
        st.stop()

    c1, c2 = st.columns([1, 2])
    with c1:
        brand_map = {b["name"]: b["id"] for b in brands}
        sel_name  = st.selectbox("Select Brand", list(brand_map.keys()))
        sel_brand = get_brand(brand_map[sel_name])
        st.markdown(
            f"<div class='brand-row'><strong>{sel_brand['name']}</strong><br>"
            f"<small>{sel_brand['industry']}</small>"
            + (f"<br><small style='color:#555'>{sel_brand['target_market']}</small>" if sel_brand.get("target_market") else "")
            + "</div>", unsafe_allow_html=True,
        )
    with c2:
        uploaded = st.file_uploader("Upload Strategy Deck (PDF or PPTX)", type=["pdf", "pptx"])

    if uploaded:
        st.markdown(f"**File:** {uploaded.name} · {uploaded.size / 1024:.1f} KB")
        if st.button("🔍 Run Analysis", type="primary", use_container_width=True):
            with st.spinner("Extracting slide content…"):
                try:
                    raw     = uploaded.read()
                    content = extract_text(raw, uploaded.name)
                    if not content.strip():
                        st.error("No readable text found. Ensure the file isn't image-only.")
                        st.stop()
                    st.success(f"Extracted {len(content.split())} words.")
                except Exception as e:
                    st.error(f"Parse error: {e}")
                    st.stop()

            with st.spinner("AI is evaluating your strategy… (15–30 seconds)"):
                try:
                    results = analyze_strategy(
                        brand_name=sel_brand["name"],
                        industry=sel_brand["industry"],
                        brand_description=sel_brand.get("description", ""),
                        target_market=sel_brand.get("target_market", ""),
                        deck_content=content,
                        deck_name=uploaded.name,
                    )
                except Exception as e:
                    st.error(f"Analysis error: {e}")
                    st.stop()

            save_analysis(sel_brand["id"], uploaded.name,
                          results["overall_score"], results["verdict"], results)
            st.success("Analysis complete! Saved to History.")
            st.session_state["last_results"] = results

    if "last_results" in st.session_state:
        st.markdown("---")
        st.subheader("Analysis Results")
        render_results(st.session_state["last_results"])


# ══════════════════════════════════════════════════════════════════════════════
# HISTORY
# ══════════════════════════════════════════════════════════════════════════════
elif page == "📋 History":
    st.title("Analysis History")
    brands = get_brands()
    f_opt = st.selectbox("Filter by Brand", ["All Brands"] + [b["name"] for b in brands])
    all_a = get_analyses() if f_opt == "All Brands" else get_analyses(next(b["id"] for b in brands if b["name"] == f_opt))

    if not all_a:
        st.info("No analyses found.")
    else:
        st.markdown(f"**{len(all_a)} result(s)**")
        for a in all_a:
            sc = score_color(a["overall_score"])
            vc = verdict_color(a["verdict"])
            with st.expander(f"📄 {a['deck_name']}  ·  {a['brand_name']}  ·  {a['overall_score']}/100"):
                results = json.loads(a["results_json"])
                mc1, mc2, mc3 = st.columns([1, 1, 2])
                with mc1:
                    st.metric("Overall Score", f"{a['overall_score']}/100")
                with mc2:
                    st.markdown(f"**Verdict**<br><span style='color:{vc};font-weight:700'>{a['verdict']}</span>", unsafe_allow_html=True)
                with mc3:
                    st.markdown(f"**Brand:** {a['brand_name']} · {a['industry']}")
                    st.markdown(f"**Date:** {a['created_at'][:16]}")
                st.markdown(results.get("executive_summary", ""))
                dims_list = [{"label": v["label"], "score": v["score"]} for v in results["dimensions"].values()]
                st.plotly_chart(bar_chart(dims_list), use_container_width=True)
                st.markdown("**Top Recommendations:**")
                pc = {"High": "#e74c3c", "Medium": "#f39c12", "Low": "#27ae60"}
                for rec in results.get("recommendations", [])[:3]:
                    c = pc.get(rec["priority"], "#666")
                    st.markdown(
                        f"<div class='rec-item' style='border-left:3px solid {c};background:{c}10'>"
                        f"<span style='color:{c};font-weight:700;font-size:.78rem'>{rec['priority'].upper()}</span> {rec['action']}</div>",
                        unsafe_allow_html=True,
                    )
                if st.button("🗑️ Delete", key=f"del_{a['id']}"):
                    delete_analysis(a["id"])
                    st.rerun()
