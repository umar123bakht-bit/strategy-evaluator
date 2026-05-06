import io, json, sqlite3, base64
import streamlit as st
import plotly.graph_objects as go
import plotly.express as px
import anthropic
from pptx import Presentation
from pptx.util import Pt
import pdfplumber
from PIL import Image

# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(page_title="Strategy Evaluator", page_icon="📊", layout="wide", initial_sidebar_state="expanded")

# ── CSS ───────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&display=swap');

html, body, [class*="css"] { font-family: 'Inter', sans-serif; }

.stApp { background: #F1F5F9; }

/* Hero */
.hero {
  background: linear-gradient(135deg, #0F2044 0%, #1B3A6B 50%, #2D5FA6 100%);
  padding: 2.5rem 3rem; border-radius: 16px; color: white;
  margin-bottom: 2rem; position: relative; overflow: hidden;
}
.hero::after {
  content: ''; position: absolute; top: -50%; right: -10%;
  width: 400px; height: 400px; border-radius: 50%;
  background: rgba(201,168,76,0.12); pointer-events: none;
}
.hero h1 { margin: 0; font-size: 2.2rem; font-weight: 800; letter-spacing: -0.5px; }
.hero p  { margin: 0.5rem 0 0; opacity: .8; font-size: 1rem; font-weight: 400; }
.hero .tag {
  display: inline-block; background: rgba(201,168,76,0.25); color: #F0CB6A;
  border: 1px solid rgba(201,168,76,0.4); border-radius: 20px;
  padding: 0.2rem 0.8rem; font-size: 0.75rem; font-weight: 600;
  margin-top: 0.8rem; letter-spacing: 0.5px;
}

/* Cards */
.card {
  background: white; border-radius: 14px; padding: 1.5rem;
  border: 1px solid #E2E8F0; box-shadow: 0 1px 3px rgba(0,0,0,.06);
  transition: box-shadow .2s, transform .2s; margin-bottom: 1rem;
}
.card:hover { box-shadow: 0 8px 25px rgba(0,0,0,.1); transform: translateY(-2px); }

/* KPI boxes */
.kpi {
  background: white; border-radius: 14px; padding: 1.4rem 1rem;
  text-align: center; border: 1px solid #E2E8F0;
  box-shadow: 0 1px 3px rgba(0,0,0,.05);
  transition: box-shadow .2s, transform .2s;
}
.kpi:hover { box-shadow: 0 6px 20px rgba(0,0,0,.1); transform: translateY(-2px); }
.kpi .val { font-size: 2.2rem; font-weight: 800; line-height: 1; }
.kpi .lbl { font-size: .75rem; color: #64748B; margin-top: .4rem; font-weight: 500; text-transform: uppercase; letter-spacing: .5px; }

/* Score badge */
.score-badge {
  display: inline-flex; align-items: center; justify-content: center;
  width: 56px; height: 56px; border-radius: 50%;
  font-size: 1.1rem; font-weight: 800; color: white;
  box-shadow: 0 4px 12px rgba(0,0,0,.2);
}

/* Verdict chips */
.verdict {
  display: inline-block; padding: .35rem 1.1rem; border-radius: 30px;
  font-weight: 700; font-size: .82rem; letter-spacing: .3px;
}

/* Dimension cards */
.dim-card {
  background: white; border-radius: 12px; border: 1px solid #E2E8F0;
  overflow: hidden; margin-bottom: .8rem;
  transition: box-shadow .2s;
}
.dim-card:hover { box-shadow: 0 6px 20px rgba(0,0,0,.08); }
.dim-header {
  padding: 1rem 1.2rem; cursor: pointer;
  display: flex; align-items: center; justify-content: space-between;
  border-left: 4px solid #2D5FA6;
}
.dim-body { padding: 0 1.2rem 1.2rem; border-top: 1px solid #F1F5F9; }

/* Progress bar */
.prog-wrap { background: #F1F5F9; border-radius: 99px; height: 8px; overflow: hidden; margin: .4rem 0; }
.prog-bar { height: 100%; border-radius: 99px; transition: width 1s ease; }

/* Framework pills */
.fw-pill {
  display: inline-block; padding: .2rem .7rem; border-radius: 99px;
  background: #EFF6FF; color: #1D4ED8; border: 1px solid #BFDBFE;
  font-size: .72rem; font-weight: 600; margin: .15rem;
}

/* Rec cards */
.rec-card {
  padding: .8rem 1rem .8rem 1.2rem; border-radius: 0 10px 10px 0;
  margin: .5rem 0; transition: transform .15s;
}
.rec-card:hover { transform: translateX(4px); }

/* Visual score ring */
.score-ring-wrap { display: flex; align-items: center; gap: 1.5rem; padding: 1rem 0; }
.score-ring {
  width: 90px; height: 90px; border-radius: 50%;
  display: flex; align-items: center; justify-content: center;
  font-size: 1.6rem; font-weight: 800; color: white;
  box-shadow: 0 6px 20px rgba(0,0,0,.2); flex-shrink: 0;
}

/* Brand row */
.brand-row {
  background: white; border-left: 4px solid #1B3A6B;
  padding: .9rem 1.1rem; border-radius: 0 10px 10px 0;
  margin-bottom: .7rem; border: 1px solid #E2E8F0;
  transition: box-shadow .2s;
}
.brand-row:hover { box-shadow: 0 4px 15px rgba(0,0,0,.08); }

/* Section title */
.section-title {
  font-size: 1.1rem; font-weight: 700; color: #0F2044;
  margin: 1.5rem 0 .8rem; padding-bottom: .4rem;
  border-bottom: 2px solid #E2E8F0;
}

/* Insight tooltip */
.insight-wrap { position: relative; display: inline-block; }
.insight-tip {
  display: none; position: absolute; z-index: 100;
  background: #0F2044; color: white; padding: .6rem .9rem;
  border-radius: 8px; font-size: .78rem; width: 200px;
  bottom: 120%; left: 50%; transform: translateX(-50%);
  box-shadow: 0 4px 15px rgba(0,0,0,.3);
}
.insight-wrap:hover .insight-tip { display: block; }

/* Stat boxes */
.stat-box {
  background: #F8FAFC; border: 1px solid #E2E8F0; border-radius: 10px;
  padding: .8rem 1rem; text-align: center;
}
.stat-box .sv { font-size: 1.4rem; font-weight: 700; }
.stat-box .sl { font-size: .72rem; color: #64748B; text-transform: uppercase; letter-spacing: .4px; margin-top: .2rem; }

/* Nav */
[data-testid="stSidebar"] { background: #0F2044 !important; }
[data-testid="stSidebar"] * { color: #CBD5E1 !important; }
[data-testid="stSidebar"] .stRadio > label { color: white !important; }

/* Divider */
.divider { border: none; border-top: 1px solid #E2E8F0; margin: 1.5rem 0; }

/* Combined score */
.combined-score {
  background: linear-gradient(135deg, #0F2044, #1B3A6B);
  color: white; border-radius: 14px; padding: 1.5rem 2rem;
  display: flex; align-items: center; gap: 2rem; margin-bottom: 1.5rem;
}
.combined-score .big { font-size: 3.5rem; font-weight: 900; line-height: 1; }
.combined-score .sub { font-size: .85rem; opacity: .75; margin-top: .3rem; }
</style>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# DATABASE
# ══════════════════════════════════════════════════════════════════════════════
DB_PATH = "strategy_data.db"

def get_conn():
    c = sqlite3.connect(DB_PATH, check_same_thread=False)
    c.row_factory = sqlite3.Row
    return c

def init_db():
    c = get_conn()
    c.execute("""CREATE TABLE IF NOT EXISTS brands (
        id INTEGER PRIMARY KEY AUTOINCREMENT, name TEXT NOT NULL UNIQUE,
        industry TEXT NOT NULL, description TEXT, target_market TEXT,
        created_at TEXT DEFAULT CURRENT_TIMESTAMP)""")
    c.execute("""CREATE TABLE IF NOT EXISTS analyses (
        id INTEGER PRIMARY KEY AUTOINCREMENT, brand_id INTEGER NOT NULL,
        deck_name TEXT NOT NULL, overall_score INTEGER, visual_score INTEGER,
        combined_score INTEGER, verdict TEXT, results_json TEXT,
        created_at TEXT DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (brand_id) REFERENCES brands(id))""")
    c.commit(); c.close()

def add_brand(name, industry, desc, market):
    c = get_conn()
    try:
        c.execute("INSERT INTO brands (name,industry,description,target_market) VALUES(?,?,?,?)", (name,industry,desc,market))
        c.commit(); return True
    except sqlite3.IntegrityError: return False
    finally: c.close()

def get_brands():
    c = get_conn(); rows = c.execute("SELECT * FROM brands ORDER BY name").fetchall(); c.close()
    return [dict(r) for r in rows]

def get_brand(bid):
    c = get_conn(); r = c.execute("SELECT * FROM brands WHERE id=?", (bid,)).fetchone(); c.close()
    return dict(r) if r else None

def update_brand(bid, name, industry, desc, market):
    c = get_conn()
    try:
        c.execute("UPDATE brands SET name=?,industry=?,description=?,target_market=? WHERE id=?", (name,industry,desc,market,bid))
        c.commit(); return True
    except sqlite3.IntegrityError: return False
    finally: c.close()

def delete_brand(bid):
    c = get_conn(); c.execute("DELETE FROM analyses WHERE brand_id=?",(bid,)); c.execute("DELETE FROM brands WHERE id=?",(bid,)); c.commit(); c.close()

def save_analysis(brand_id, deck_name, overall_score, visual_score, combined_score, verdict, results_json):
    c = get_conn()
    cur = c.execute("INSERT INTO analyses (brand_id,deck_name,overall_score,visual_score,combined_score,verdict,results_json) VALUES(?,?,?,?,?,?,?)",
        (brand_id, deck_name, overall_score, visual_score, combined_score, verdict, json.dumps(results_json)))
    aid = cur.lastrowid; c.commit(); c.close(); return aid

def get_analyses(brand_id=None):
    c = get_conn()
    q = """SELECT a.*,b.name as brand_name,b.industry FROM analyses a JOIN brands b ON a.brand_id=b.id"""
    rows = c.execute(q + (" WHERE a.brand_id=? ORDER BY a.created_at DESC" if brand_id else " ORDER BY a.created_at DESC"),
        (brand_id,) if brand_id else ()).fetchall()
    c.close(); return [dict(r) for r in rows]

def delete_analysis(aid):
    c = get_conn(); c.execute("DELETE FROM analyses WHERE id=?",(aid,)); c.commit(); c.close()


# ══════════════════════════════════════════════════════════════════════════════
# FILE PARSING
# ══════════════════════════════════════════════════════════════════════════════
def extract_text_and_visual(file_bytes, filename):
    """Returns (text_content, visual_metadata, image_samples)"""
    name = filename.lower()

    if name.endswith(".pptx"):
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts, slides_info, image_samples = [], [], []

        for i, slide in enumerate(prs.slides, 1):
            words, imgs, shapes = 0, 0, len(slide.shapes)
            has_chart = has_table = False
            slide_texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                    words += len(shape.text.split())
                if shape.shape_type == 13:
                    imgs += 1
                    if len(image_samples) < 6:
                        try: image_samples.append(shape.image.blob)
                        except: pass
                if shape.shape_type == 3: has_chart = True
                if shape.shape_type == 19: has_table = True

            slides_info.append({"slide":i,"words":words,"images":imgs,"shapes":shapes,"chart":has_chart,"table":has_table})
            if slide_texts: text_parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))

        total = len(slides_info) or 1
        visual_meta = {
            "file_type": "PPTX Presentation",
            "total_slides": len(slides_info),
            "total_images": sum(s["images"] for s in slides_info),
            "total_charts": sum(1 for s in slides_info if s["chart"]),
            "total_tables": sum(1 for s in slides_info if s["table"]),
            "avg_words_per_slide": round(sum(s["words"] for s in slides_info) / total, 1),
            "text_heavy_slides": sum(1 for s in slides_info if s["words"] > 80),
            "visual_slides": sum(1 for s in slides_info if s["images"] > 0 or s["chart"]),
            "max_words_slide": max((s["words"] for s in slides_info), default=0),
            "image_coverage_pct": round(sum(1 for s in slides_info if s["images"] > 0) / total * 100),
        }
        return "\n\n".join(text_parts), visual_meta, image_samples

    elif name.endswith(".pdf"):
        text_parts, pages_info, image_samples = [], [], []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                imgs = page.images or []
                words = len(text.split())
                pages_info.append({"page":i,"words":words,"images":len(imgs)})
                if text.strip(): text_parts.append(f"[Page {i}]\n{text.strip()}")

        total = len(pages_info) or 1
        visual_meta = {
            "file_type": "PDF Document",
            "total_pages": len(pages_info),
            "total_images": sum(p["images"] for p in pages_info),
            "avg_words_per_page": round(sum(p["words"] for p in pages_info) / total, 1),
            "text_heavy_pages": sum(1 for p in pages_info if p["words"] > 200),
            "visual_pages": sum(1 for p in pages_info if p["images"] > 0),
            "max_words_page": max((p["words"] for p in pages_info), default=0),
            "image_coverage_pct": round(sum(1 for p in pages_info if p["images"] > 0) / total * 100),
        }
        return "\n\n".join(text_parts), visual_meta, image_samples

    raise ValueError(f"Unsupported: {filename}")


# ══════════════════════════════════════════════════════════════════════════════
# AI — STRATEGY ANALYSIS (McKinsey-level)
# ══════════════════════════════════════════════════════════════════════════════
def analyze_strategy(brand_name, industry, brand_description, target_market, deck_content, deck_name):
    client = anthropic.Anthropic(api_key=st.secrets["ANTHROPIC_API_KEY"])
    prompt = f"""You are a Senior Partner at McKinsey & Company with 25 years of experience leading strategy engagements for Fortune 500 companies globally. Your assessments are known for intellectual rigor, honesty, and transformative impact.

BRAND: {brand_name} | INDUSTRY: {industry}
DESCRIPTION: {brand_description or "Not provided"}
TARGET MARKET: {target_market or "Not provided"}
DOCUMENT: {deck_name}

STRATEGY DOCUMENT:
{deck_content[:14000]}

━━━ ANALYTICAL FRAMEWORKS TO APPLY ━━━
1. "Where to Play / How to Win" (Lafley & Martin) — Are strategic choices explicit?
2. Porter's Five Forces — Competitive dynamics and structural advantage
3. Jobs-to-be-Done — Does the value proposition solve real, significant problems?
4. BCG Matrix — Market attractiveness vs competitive strength
5. McKinsey 7S — Strategy, Structure, Systems, Staff, Skills, Style, Shared Values
6. Ansoff Matrix — Growth vector clarity (penetration / development / diversification)
7. MECE Principle — Mutually Exclusive, Collectively Exhaustive thinking
8. Blue Ocean Canvas — Value innovation vs competing in red ocean

━━━ SCORING RUBRIC (apply with full McKinsey rigor — DO NOT inflate) ━━━
90–100: World-class. Sets industry benchmark. Extremely rare.
80–89: Excellent. Strong strategic thinking, minor refinements needed.
70–79: Good. Solid fundamentals with meaningful gaps.
60–69: Average. Conventional thinking, lacks differentiation.
50–59: Below average. Significant strategic weaknesses present.
Below 50: Poor. Fundamental flaws requiring major rethinking.

CRITICAL: A score above 75 must be genuinely earned. Challenge every assumption. Flag what is MISSING as much as what is present. Apply the "So What? Now What?" test to every finding.

Return ONLY valid JSON — no markdown, no extra text:
{{
  "overall_score": <integer 0-100>,
  "executive_summary": "<3 sentences: what the strategy does, its biggest strength, its most critical gap>",
  "verdict": "<Highly Effective | Effective | Needs Improvement | Ineffective>",
  "frameworks_applied": ["<framework1>", "<framework2>", "<framework3>"],
  "dimensions": {{
    "strategic_clarity": {{"score": <0-100>, "label": "Strategic Clarity & Coherence", "feedback": "<specific McKinsey-level feedback>", "so_what": "<implication>", "strengths": ["<s>","<s>"], "improvements": ["<i>","<i>"]}},
    "market_opportunity": {{"score": <0-100>, "label": "Market Opportunity & Timing", "feedback": "<specific>", "so_what": "<implication>", "strengths": ["<s>"], "improvements": ["<i>","<i>"]}},
    "competitive_advantage": {{"score": <0-100>, "label": "Competitive Advantage & Moat", "feedback": "<specific>", "so_what": "<implication>", "strengths": ["<s>"], "improvements": ["<i>","<i>"]}},
    "customer_value": {{"score": <0-100>, "label": "Customer & Value Proposition", "feedback": "<specific>", "so_what": "<implication>", "strengths": ["<s>"], "improvements": ["<i>","<i>"]}},
    "financial_viability": {{"score": <0-100>, "label": "Financial Viability & ROI", "feedback": "<specific>", "so_what": "<implication>", "strengths": ["<s>"], "improvements": ["<i>","<i>"]}},
    "execution_readiness": {{"score": <0-100>, "label": "Execution Readiness", "feedback": "<specific>", "so_what": "<implication>", "strengths": ["<s>"], "improvements": ["<i>","<i>"]}},
    "risk_resilience": {{"score": <0-100>, "label": "Risk & Resilience", "feedback": "<specific>", "so_what": "<implication>", "strengths": ["<s>"], "improvements": ["<i>","<i>"]}},
    "measurement": {{"score": <0-100>, "label": "Measurement & Accountability", "feedback": "<specific>", "so_what": "<implication>", "strengths": ["<s>"], "improvements": ["<i>","<i>"]}}
  }},
  "top_strengths": ["<strength 1>", "<strength 2>", "<strength 3>"],
  "critical_gaps": ["<gap 1>", "<gap 2>", "<gap 3>"],
  "recommendations": [
    {{"priority": "Quick Win", "timeline": "0-30 days", "action": "<specific step>", "impact": "<expected outcome>"}},
    {{"priority": "Quick Win", "timeline": "0-30 days", "action": "<specific step>", "impact": "<expected outcome>"}},
    {{"priority": "Medium Term", "timeline": "1-3 months", "action": "<specific step>", "impact": "<expected outcome>"}},
    {{"priority": "Medium Term", "timeline": "1-3 months", "action": "<specific step>", "impact": "<expected outcome>"}},
    {{"priority": "Strategic", "timeline": "3-6 months", "action": "<specific step>", "impact": "<expected outcome>"}}
  ]
}}"""
    msg = client.messages.create(model="claude-opus-4-5", max_tokens=4096, messages=[{"role":"user","content":prompt}])
    text = msg.content[0].text.strip()
    if text.startswith("```"):
        lines = text.split("\n"); text = "\n".join(lines[1:-1] if lines[-1].strip()=="```" else lines[1:])
    return json.loads(text)


# ══════════════════════════════════════════════════════════════════════════════
# AI — VISUAL APPEAL ANALYSIS
# ══════════════════════════════════════════════════════════════════════════════
def analyze_visual_appeal(visual_meta, image_samples, brand_name, industry, deck_name):
    client = anthropic.Anthropic(api_key=st.secrets["ANTHROPIC_API_KEY"])

    meta_str = json.dumps(visual_meta, indent=2)
    is_pptx = visual_meta.get("file_type","").startswith("PPTX")
    slide_word = "slide" if is_pptx else "page"
    total = visual_meta.get("total_slides", visual_meta.get("total_pages", 0))
    avg_words = visual_meta.get("avg_words_per_slide", visual_meta.get("avg_words_per_page", 0))
    text_heavy = visual_meta.get("text_heavy_slides", visual_meta.get("text_heavy_pages", 0))
    visual_count = visual_meta.get("visual_slides", visual_meta.get("visual_pages", 0))

    prompt_text = f"""You are a world-class visual communications director and presentation design expert. You have reviewed thousands of strategy decks for global consulting firms, investment banks, and Fortune 500 companies.

Evaluate the VISUAL DESIGN and COMMUNICATION QUALITY of this deck — not its strategic content.

BRAND: {brand_name} | INDUSTRY: {industry}
DECK: {deck_name}

STRUCTURAL METRICS:
{meta_str}

KEY DESIGN INDICATORS:
- Average words per {slide_word}: {avg_words} {"⚠️ HIGH — likely text-heavy slides" if avg_words > 80 else "✅ Good balance" if avg_words < 60 else "⚡ Borderline"}
- Text-heavy {slide_word}s (>80 words): {text_heavy} of {total} ({round(text_heavy/max(total,1)*100)}%)
- {slide_word.title()}s with visuals: {visual_count} of {total} ({round(visual_count/max(total,1)*100)}%)
- Charts/graphs: {visual_meta.get("total_charts", "N/A")} | Tables: {visual_meta.get("total_tables", "N/A")}

SCORING (be strict and design-critical):
90-100: Exceptional. McKinsey/BCG tier. Every slide communicates visually. Memorable.
75-89: Professional. Strong visual design with minor issues.
60-74: Average. Functional but over-reliant on text, lacks visual impact.
45-59: Below average. Text-dumps, weak design hierarchy.
0-44: Poor. No visual strategy. Walls of text. Undermines the strategy's credibility.

IMPORTANT: Decks where >40% of slides are text-heavy should score below 65 on visual balance. Decks with <30% image coverage should score below 60 on visual storytelling.

Return ONLY valid JSON:
{{
  "visual_score": <integer 0-100>,
  "visual_verdict": "<Visually Exceptional | Visually Strong | Visually Average | Needs Design Work | Poor Visual Design>",
  "visual_summary": "<2-3 sentences — be specific about what works and what doesn't visually>",
  "dimensions": {{
    "text_visual_balance": {{"score": <0-100>, "label": "Text-Visual Balance", "feedback": "<specific>", "finding": "<key metric or observation>"}},
    "design_consistency": {{"score": <0-100>, "label": "Design Consistency", "feedback": "<specific>", "finding": "<key observation>"}},
    "data_visualization": {{"score": <0-100>, "label": "Data Visualization", "feedback": "<specific>", "finding": "<key observation>"}},
    "layout_whitespace": {{"score": <0-100>, "label": "Layout & White Space", "feedback": "<specific>", "finding": "<key observation>"}},
    "visual_storytelling": {{"score": <0-100>, "label": "Visual Storytelling", "feedback": "<specific>", "finding": "<key observation>"}},
    "professional_polish": {{"score": <0-100>, "label": "Professional Polish", "feedback": "<specific>", "finding": "<key observation>"}}
  }},
  "design_strengths": ["<specific strength>", "<specific strength>"],
  "design_issues": ["<specific issue>", "<specific issue>", "<specific issue>"],
  "design_recommendations": [
    {{"priority": "High", "action": "<specific design fix>", "why": "<impact on audience>"}},
    {{"priority": "High", "action": "<specific design fix>", "why": "<impact>"}},
    {{"priority": "Medium", "action": "<specific design improvement>", "why": "<impact>"}}
  ],
  "slide_stats": {{
    "total_slides": {total},
    "avg_words_per_slide": {avg_words},
    "pct_visual_slides": {round(visual_count/max(total,1)*100)},
    "pct_text_heavy": {round(text_heavy/max(total,1)*100)}
  }}
}}"""

    # Build multimodal content — add extracted images if available
    def resize_image(blob, max_px=1500):
        """Resize image blob so neither dimension exceeds max_px."""
        try:
            img = Image.open(io.BytesIO(blob)).convert("RGB")
            img.thumbnail((max_px, max_px), Image.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=75)
            return buf.getvalue()
        except:
            return None

    content = []
    for blob in image_samples[:4]:
        if blob:
            resized = resize_image(blob)
            if resized:
                content.append({"type":"image","source":{"type":"base64","media_type":"image/jpeg","data":base64.b64encode(resized).decode()}})
    content.append({"type":"text","text":prompt_text})

    msg = client.messages.create(model="claude-opus-4-5", max_tokens=2048, messages=[{"role":"user","content":content}])
    text = msg.content[0].text.strip()
    if text.startswith("```"):
        lines = text.split("\n"); text = "\n".join(lines[1:-1] if lines[-1].strip()=="```" else lines[1:])
    return json.loads(text)


# ══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════════════
def sc(s):
    if s >= 80: return "#059669"
    if s >= 65: return "#2563EB"
    if s >= 50: return "#D97706"
    return "#DC2626"

def vc(v):
    return {"Highly Effective":"#059669","Effective":"#2563EB","Needs Improvement":"#D97706","Ineffective":"#DC2626",
            "Visually Exceptional":"#059669","Visually Strong":"#2563EB","Visually Average":"#D97706",
            "Needs Design Work":"#D97706","Poor Visual Design":"#DC2626"}.get(v,"#64748B")

def score_label(s):
    if s >= 80: return "Excellent"
    if s >= 65: return "Good"
    if s >= 50: return "Average"
    return "Needs Work"

def radar(dims, color="#2563EB"):
    labels = [d["label"] for d in dims] + [dims[0]["label"]]
    scores = [d["score"] for d in dims] + [dims[0]["score"]]
    fig = go.Figure(go.Scatterpolar(r=scores, theta=labels, fill="toself",
        fillcolor=f"rgba(37,99,235,.12)", line=dict(color=color, width=2.5),
        hovertemplate="<b>%{theta}</b><br>Score: %{r}/100<extra></extra>"))
    fig.update_layout(polar=dict(radialaxis=dict(visible=True,range=[0,100],tickfont=dict(size=10)),
        angularaxis=dict(tickfont=dict(size=11))), showlegend=False, height=380,
        margin=dict(l=55,r=55,t=30,b=30), paper_bgcolor="rgba(0,0,0,0)")
    return fig

def hbar(dims):
    fig = go.Figure(go.Bar(x=[d["score"] for d in dims], y=[d["label"] for d in dims],
        orientation="h", marker_color=[sc(d["score"]) for d in dims],
        text=[f"  {d['score']}/100  {score_label(d['score'])}" for d in dims],
        textposition="outside",
        hovertemplate="<b>%{y}</b><br>Score: %{x}/100<extra></extra>"))
    fig.update_layout(xaxis=dict(range=[0,120],showgrid=True,gridcolor="#F1F5F9"),
        yaxis=dict(autorange="reversed"), height=320, plot_bgcolor="white",
        paper_bgcolor="rgba(0,0,0,0)", margin=dict(l=10,r=80,t=10,b=10))
    return fig

def gauge(score, title=""):
    fig = go.Figure(go.Indicator(mode="gauge+number", value=score,
        title={"text":title,"font":{"size":13}},
        gauge={"axis":{"range":[0,100]},"bar":{"color":sc(score),"thickness":.25},
               "steps":[{"range":[0,50],"color":"#FEF2F2"},{"range":[50,65],"color":"#FFFBEB"},
                        {"range":[65,80],"color":"#EFF6FF"},{"range":[80,100],"color":"#ECFDF5"}],
               "threshold":{"line":{"color":"#1E293B","width":2},"thickness":.75,"value":score}}))
    fig.update_layout(height=200, margin=dict(l=20,r=20,t=30,b=10), paper_bgcolor="rgba(0,0,0,0)")
    return fig


# ══════════════════════════════════════════════════════════════════════════════
# RENDER FULL RESULTS
# ══════════════════════════════════════════════════════════════════════════════
def render_full_results(results):
    strat = results["strategy"]
    visual = results["visual"]
    combined = results["combined_score"]

    # ── Combined header ──
    vc_color = vc(strat["verdict"])
    st.markdown(f"""
    <div class="combined-score">
      <div>
        <div class="big" style="color:#C9A84C">{combined}</div>
        <div style="font-size:.8rem;opacity:.7;margin-top:.3rem">COMBINED SCORE /100</div>
      </div>
      <div style="flex:1">
        <div style="font-size:1.3rem;font-weight:700">{strat['verdict']}</div>
        <div style="opacity:.75;font-size:.88rem;margin-top:.3rem">{strat['executive_summary']}</div>
        <div style="margin-top:.8rem">
          {''.join(f'<span class="fw-pill">{f}</span>' for f in strat.get("frameworks_applied",[])[:5])}
        </div>
      </div>
      <div style="text-align:center">
        <div style="font-size:.75rem;opacity:.6;text-transform:uppercase;letter-spacing:.5px">Strategy</div>
        <div style="font-size:2.2rem;font-weight:800;color:{'#C9A84C' if strat['overall_score']>=65 else '#F87171'}">{strat['overall_score']}</div>
        <div style="font-size:.75rem;opacity:.6;margin-top:.5rem;text-transform:uppercase;letter-spacing:.5px">Visual</div>
        <div style="font-size:2.2rem;font-weight:800;color:{'#C9A84C' if visual['visual_score']>=65 else '#F87171'}">{visual['visual_score']}</div>
      </div>
    </div>
    """, unsafe_allow_html=True)

    # ── Tabs ──
    t1, t2, t3, t4 = st.tabs(["📊 Strategy Analysis", "🎨 Visual Design", "💡 Recommendations", "📈 Charts"])

    # ── TAB 1: Strategy ──
    with t1:
        c1, c2 = st.columns(2)
        with c1:
            st.markdown("<div class='section-title'>💪 Top Strengths</div>", unsafe_allow_html=True)
            for s in strat.get("top_strengths", []):
                st.markdown(f"""<div class="card" style="padding:.8rem 1rem;border-left:3px solid #059669">
                ✅ &nbsp;{s}</div>""", unsafe_allow_html=True)
        with c2:
            st.markdown("<div class='section-title'>⚠️ Critical Gaps</div>", unsafe_allow_html=True)
            for g in strat.get("critical_gaps", []):
                st.markdown(f"""<div class="card" style="padding:.8rem 1rem;border-left:3px solid #DC2626">
                🔴 &nbsp;{g}</div>""", unsafe_allow_html=True)

        st.markdown("<div class='section-title'>🔍 Dimension Deep-Dives — Click any to expand</div>", unsafe_allow_html=True)

        for key, dim in strat["dimensions"].items():
            color = sc(dim["score"])
            exp_key = f"strat_dim_{key}"
            with st.expander(f"{'🟢' if dim['score']>=80 else '🟡' if dim['score']>=60 else '🔴'}  {dim['label']}  —  {dim['score']}/100  ·  {score_label(dim['score'])}"):
                st.markdown(f"""
                <div style="margin-bottom:.8rem">
                  <div style="background:#F8FAFC;border-radius:8px;padding:.8rem 1rem;border-left:4px solid {color}">
                    <div style="font-size:.82rem;color:#64748B;font-weight:600;text-transform:uppercase;letter-spacing:.4px">Analysis</div>
                    <div style="margin-top:.3rem">{dim['feedback']}</div>
                  </div>
                </div>
                """, unsafe_allow_html=True)

                if dim.get("so_what"):
                    st.markdown(f"""<div style="background:#EFF6FF;border:1px solid #BFDBFE;border-radius:8px;padding:.7rem 1rem;margin-bottom:.8rem">
                    <span style="color:#1D4ED8;font-weight:700;font-size:.8rem">💡 SO WHAT?</span><br>
                    <span style="font-size:.9rem">{dim['so_what']}</span></div>""", unsafe_allow_html=True)

                prog = int(dim['score'])
                st.markdown(f"""<div style="margin:.5rem 0 1rem">
                  <div style="display:flex;justify-content:space-between;font-size:.8rem;color:#64748B;margin-bottom:.3rem">
                    <span>Score</span><span style="font-weight:700;color:{color}">{prog}/100</span>
                  </div>
                  <div class="prog-wrap"><div class="prog-bar" style="width:{prog}%;background:{color}"></div></div>
                </div>""", unsafe_allow_html=True)

                c1, c2 = st.columns(2)
                with c1:
                    if dim.get("strengths"):
                        st.markdown("**✅ Strengths**")
                        for s in dim["strengths"]: st.markdown(f"- {s}")
                with c2:
                    if dim.get("improvements"):
                        st.markdown("**🔧 Improvements**")
                        for i in dim["improvements"]: st.markdown(f"- {i}")

    # ── TAB 2: Visual ──
    with t2:
        vv_color = vc(visual["visual_verdict"])
        stats = visual.get("slide_stats", {})

        st.markdown(f"""
        <div class="card" style="border-left:4px solid {vv_color};margin-bottom:1.2rem">
          <div style="display:flex;align-items:center;gap:1.2rem">
            <div class="score-ring" style="background:{vv_color}">{visual['visual_score']}</div>
            <div>
              <div style="font-size:1.2rem;font-weight:700">{visual['visual_verdict']}</div>
              <div style="color:#64748B;margin-top:.3rem;font-size:.9rem">{visual['visual_summary']}</div>
            </div>
          </div>
        </div>
        """, unsafe_allow_html=True)

        # Stats row
        s1, s2, s3, s4 = st.columns(4)
        stats_data = [
            (stats.get("total_slides","—"), "Total Slides"),
            (f"{stats.get('avg_words_per_slide','—')}", "Avg Words/Slide"),
            (f"{stats.get('pct_visual_slides','—')}%", "Slides w/ Visuals"),
            (f"{stats.get('pct_text_heavy','—')}%", "Text-Heavy Slides"),
        ]
        for col, (val, lbl) in zip([s1,s2,s3,s4], stats_data):
            col.markdown(f"""<div class="stat-box"><div class="sv">{val}</div><div class="sl">{lbl}</div></div>""", unsafe_allow_html=True)

        st.markdown("<div class='section-title'>🎨 Visual Dimension Scores — Click to expand</div>", unsafe_allow_html=True)

        for key, dim in visual["dimensions"].items():
            color = sc(dim["score"])
            with st.expander(f"{'🟢' if dim['score']>=80 else '🟡' if dim['score']>=60 else '🔴'}  {dim['label']}  —  {dim['score']}/100"):
                if dim.get("finding"):
                    st.markdown(f"""<div style="background:#FFF7ED;border:1px solid #FED7AA;border-radius:8px;padding:.7rem 1rem;margin-bottom:.8rem">
                    <span style="color:#C2410C;font-weight:700;font-size:.8rem">📌 KEY FINDING</span><br>
                    <span style="font-size:.9rem">{dim['finding']}</span></div>""", unsafe_allow_html=True)
                st.markdown(dim["feedback"])
                prog = int(dim["score"])
                st.markdown(f"""<div class="prog-wrap" style="margin-top:.5rem"><div class="prog-bar" style="width:{prog}%;background:{color}"></div></div>""", unsafe_allow_html=True)

        c1, c2 = st.columns(2)
        with c1:
            st.markdown("<div class='section-title'>✅ Design Strengths</div>", unsafe_allow_html=True)
            for s in visual.get("design_strengths", []):
                st.markdown(f"""<div class="card" style="padding:.7rem 1rem;border-left:3px solid #059669;font-size:.9rem">✅ {s}</div>""", unsafe_allow_html=True)
        with c2:
            st.markdown("<div class='section-title'>🚨 Design Issues</div>", unsafe_allow_html=True)
            for i in visual.get("design_issues", []):
                st.markdown(f"""<div class="card" style="padding:.7rem 1rem;border-left:3px solid #DC2626;font-size:.9rem">❌ {i}</div>""", unsafe_allow_html=True)

        st.markdown("<div class='section-title'>🛠️ Design Recommendations</div>", unsafe_allow_html=True)
        pc = {"High":"#DC2626","Medium":"#D97706","Low":"#059669"}
        for rec in visual.get("design_recommendations", []):
            c = pc.get(rec["priority"],"#64748B")
            st.markdown(f"""<div class="rec-card" style="border-left:4px solid {c};background:{c}10">
            <span style="color:{c};font-weight:700;font-size:.78rem">{rec['priority'].upper()}</span>
            <span style="font-size:.9rem;margin-left:.5rem">{rec['action']}</span>
            {f'<br><span style="font-size:.8rem;color:#64748B;margin-left:.5rem">→ {rec["why"]}</span>' if rec.get("why") else ""}
            </div>""", unsafe_allow_html=True)

    # ── TAB 3: Recommendations ──
    with t3:
        st.markdown("<div class='section-title'>🎯 Strategic Recommendations</div>", unsafe_allow_html=True)
        pc = {"Quick Win":"#059669","Medium Term":"#2563EB","Strategic":"#7C3AED"}
        for rec in strat.get("recommendations", []):
            c = pc.get(rec["priority"],"#64748B")
            st.markdown(f"""<div class="card" style="border-left:4px solid {c};padding:1rem 1.2rem">
            <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:.4rem">
              <span style="background:{c}18;color:{c};border:1px solid {c}40;border-radius:99px;padding:.2rem .8rem;font-size:.78rem;font-weight:700">{rec['priority']}</span>
              <span style="color:#64748B;font-size:.8rem">⏱ {rec.get('timeline','')}</span>
            </div>
            <div style="font-weight:600;font-size:.95rem">{rec['action']}</div>
            {f'<div style="color:#64748B;font-size:.85rem;margin-top:.3rem">→ {rec["impact"]}</div>' if rec.get("impact") else ""}
            </div>""", unsafe_allow_html=True)

        st.markdown("<div class='section-title'>🎨 Design Recommendations</div>", unsafe_allow_html=True)
        pc2 = {"High":"#DC2626","Medium":"#D97706","Low":"#059669"}
        for rec in visual.get("design_recommendations", []):
            c = pc2.get(rec["priority"],"#64748B")
            st.markdown(f"""<div class="card" style="border-left:4px solid {c};padding:1rem 1.2rem">
            <div style="display:flex;align-items:center;gap:.6rem;margin-bottom:.3rem">
              <span style="background:{c}18;color:{c};border:1px solid {c}40;border-radius:99px;padding:.2rem .8rem;font-size:.78rem;font-weight:700">{rec['priority']} Priority</span>
            </div>
            <div style="font-weight:600;font-size:.95rem">{rec['action']}</div>
            {f'<div style="color:#64748B;font-size:.85rem;margin-top:.3rem">→ {rec["why"]}</div>' if rec.get("why") else ""}
            </div>""", unsafe_allow_html=True)

    # ── TAB 4: Charts ──
    with t4:
        c1, c2 = st.columns(2)
        strat_dims = [{"label":v["label"],"score":v["score"]} for v in strat["dimensions"].values()]
        vis_dims   = [{"label":v["label"],"score":v["score"]} for v in visual["dimensions"].values()]

        with c1:
            st.subheader("Strategy Radar")
            st.plotly_chart(radar(strat_dims, "#2563EB"), use_container_width=True)
            st.subheader("Strategy Scores")
            st.plotly_chart(hbar(strat_dims), use_container_width=True)

        with c2:
            st.subheader("Visual Design Radar")
            st.plotly_chart(radar(vis_dims, "#7C3AED"), use_container_width=True)
            st.subheader("Visual Scores")
            st.plotly_chart(hbar(vis_dims), use_container_width=True)

        # Combined gauge row
        st.markdown("---")
        g1, g2, g3 = st.columns(3)
        with g1: st.plotly_chart(gauge(strat["overall_score"], "Strategy Score"), use_container_width=True)
        with g2: st.plotly_chart(gauge(visual["visual_score"], "Visual Score"), use_container_width=True)
        with g3: st.plotly_chart(gauge(combined, "Combined Score"), use_container_width=True)


# ══════════════════════════════════════════════════════════════════════════════
# APP INIT
# ══════════════════════════════════════════════════════════════════════════════
init_db()
INDUSTRIES = ["Technology","Retail & E-commerce","Food & Beverage","Fashion & Apparel",
    "Healthcare","Finance & Banking","Real Estate","Education","Entertainment & Media",
    "Automotive","Travel & Hospitality","Beauty & Personal Care","Sports & Fitness","Manufacturing","Other"]

# ── Sidebar ──
with st.sidebar:
    st.markdown("""<div style="padding:.5rem 0 1rem">
    <div style="font-size:1.3rem;font-weight:800;color:white;letter-spacing:-.3px">📊 Strategy</div>
    <div style="font-size:1.3rem;font-weight:800;color:#C9A84C;letter-spacing:-.3px">Evaluator</div>
    <div style="font-size:.75rem;color:#94A3B8;margin-top:.3rem">McKinsey-Level Analysis</div>
    </div>""", unsafe_allow_html=True)
    st.markdown("---")
    page = st.radio("", ["🏠 Dashboard","🏢 Brands","🔍 Analyze","📋 History"], label_visibility="collapsed")
    st.markdown("---")
    _b = get_brands(); _a = get_analyses()
    st.markdown(f"<div style='color:#94A3B8;font-size:.8rem'><b style='color:white'>{len(_b)}</b> brands · <b style='color:white'>{len(_a)}</b> analyses</div>", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
if page == "🏠 Dashboard":
    st.markdown("""<div class="hero">
    <div class="tag">MCKINSEY-LEVEL ANALYSIS</div>
    <h1>Strategy Evaluator</h1>
    <p>AI-powered strategy & visual design assessment for your brands</p>
    </div>""", unsafe_allow_html=True)

    avg_s = sum(a["overall_score"] for a in _a) / len(_a) if _a else 0
    avg_v = sum((a.get("visual_score") or 0) for a in _a) / len(_a) if _a else 0
    eff   = sum(1 for a in _a if a["verdict"] in ("Highly Effective","Effective"))

    c1,c2,c3,c4,c5 = st.columns(5)
    for col,val,lbl,color in [
        (c1, len(_b),        "Brands",           "#1B3A6B"),
        (c2, len(_a),        "Analyses Run",      "#1B3A6B"),
        (c3, f"{avg_s:.0f}", "Avg Strategy Score",sc(avg_s) if _a else "#94A3B8"),
        (c4, f"{avg_v:.0f}", "Avg Visual Score",  sc(avg_v) if _a else "#94A3B8"),
        (c5, eff,            "Effective",         "#059669"),
    ]:
        col.markdown(f"""<div class="kpi"><div class="val" style="color:{color}">{val}</div>
        <div class="lbl">{lbl}</div></div>""", unsafe_allow_html=True)

    if not _a:
        st.markdown("<br>", unsafe_allow_html=True)
        st.info("No analyses yet. Go to **Brands** to add your brands, then **Analyze** to run your first evaluation.")
    else:
        st.markdown("<div class='section-title'>Recent Analyses</div>", unsafe_allow_html=True)
        cl, cr = st.columns([3,1])
        with cl:
            for a in _a[:6]:
                s_c = sc(a["overall_score"]); v_c = vc(a["verdict"])
                vs = a.get("visual_score") or "—"
                cs = a.get("combined_score") or a["overall_score"]
                st.markdown(f"""<div class="brand-row">
                <div style="display:flex;justify-content:space-between;align-items:center">
                  <div><strong>{a['deck_name']}</strong> · <em style="color:#64748B">{a['brand_name']}</em></div>
                  <div style="display:flex;gap:.8rem;align-items:center">
                    <span style="font-size:.78rem;color:#64748B">Strategy <b style="color:{s_c}">{a['overall_score']}</b></span>
                    <span style="font-size:.78rem;color:#64748B">Visual <b style="color:{sc(vs) if isinstance(vs,int) else '#94A3B8'}">{vs}</b></span>
                    <span style="font-size:.82rem;font-weight:700;color:{vc(a['verdict'])}">{a['verdict']}</span>
                  </div>
                </div>
                <div style="font-size:.78rem;color:#94A3B8;margin-top:.3rem">{a['industry']} · {a['created_at'][:16]}</div>
                </div>""", unsafe_allow_html=True)

        with cr:
            st.markdown("<div style='font-weight:700;color:#0F2044;margin-bottom:.6rem'>Score Distribution</div>", unsafe_allow_html=True)
            scores = [a["overall_score"] for a in _a]
            buckets = {"0–49":0,"50–64":0,"65–79":0,"80–100":0}
            for s in scores:
                if s<50: buckets["0–49"]+=1
                elif s<65: buckets["50–64"]+=1
                elif s<80: buckets["65–79"]+=1
                else: buckets["80–100"]+=1
            fig = px.pie(values=list(buckets.values()), names=list(buckets.keys()),
                color_discrete_sequence=["#DC2626","#D97706","#2563EB","#059669"], hole=.4)
            fig.update_layout(height=220, margin=dict(l=0,r=0,t=10,b=10), paper_bgcolor="rgba(0,0,0,0)", showlegend=True, legend=dict(font=dict(size=11)))
            st.plotly_chart(fig, use_container_width=True)

    if _b:
        st.markdown("<div class='section-title'>Brand Overview</div>", unsafe_allow_html=True)
        for brand in _b:
            ba = [a for a in _a if a["brand_name"]==brand["name"]]
            avg = sum(a["overall_score"] for a in ba)/len(ba) if ba else None
            ca,cb,cc,cd = st.columns([3,1,1,1])
            ca.markdown(f"**{brand['name']}** · <span style='color:#64748B'>{brand['industry']}</span>", unsafe_allow_html=True)
            cb.markdown(f"**{len(ba)}** analyses")
            if avg: cc.markdown(f"Strategy avg: <b style='color:{sc(avg)}'>{avg:.0f}</b>", unsafe_allow_html=True)
            else: cc.markdown("No data yet")


# ══════════════════════════════════════════════════════════════════════════════
# BRANDS
# ══════════════════════════════════════════════════════════════════════════════
elif page == "🏢 Brands":
    st.markdown("<div class='hero'><h1>Brand Management</h1><p>Register and manage your brands for strategy evaluation</p></div>", unsafe_allow_html=True)
    tab_add, tab_all = st.tabs(["➕ Add Brand","📋 All Brands"])

    with tab_add:
        with st.form("add_brand"):
            c1,c2 = st.columns(2)
            with c1:
                bn = st.text_input("Brand Name *", placeholder="e.g., Acme Corp")
                bi = st.selectbox("Industry *", INDUSTRIES)
            with c2:
                bm = st.text_input("Target Market", placeholder="e.g., Young professionals 25–35")
                bd = st.text_area("Brand Description", placeholder="Values, positioning, differentiators…", height=100)
            if st.form_submit_button("Add Brand ➕", type="primary"):
                if not bn.strip(): st.error("Brand name required.")
                elif add_brand(bn.strip(), bi, bd, bm): st.success(f"✅ '{bn}' added!"); st.rerun()
                else: st.error(f"'{bn}' already exists.")

    with tab_all:
        brands = get_brands()
        if not brands: st.info("No brands yet.")
        for brand in brands:
            with st.expander(f"🏢  {brand['name']}  —  {brand['industry']}"):
                ekey = f"edit_{brand['id']}"
                if st.session_state.get(ekey):
                    with st.form(f"ef_{brand['id']}"):
                        ec1,ec2 = st.columns(2)
                        with ec1:
                            nn = st.text_input("Name", value=brand["name"])
                            ni = st.selectbox("Industry", INDUSTRIES, index=INDUSTRIES.index(brand["industry"]) if brand["industry"] in INDUSTRIES else 0)
                        with ec2:
                            nm = st.text_input("Target Market", value=brand.get("target_market") or "")
                            nd = st.text_area("Description", value=brand.get("description") or "", height=80)
                        s1,s2 = st.columns(2)
                        if s1.form_submit_button("Save ✅", type="primary"): update_brand(brand["id"],nn,ni,nd,nm); st.session_state[ekey]=False; st.rerun()
                        if s2.form_submit_button("Cancel"): st.session_state[ekey]=False; st.rerun()
                else:
                    dc1,dc2 = st.columns([3,1])
                    with dc1:
                        st.markdown(f"**Industry:** {brand['industry']}")
                        if brand.get("target_market"): st.markdown(f"**Target Market:** {brand['target_market']}")
                        if brand.get("description"): st.markdown(f"**Description:** {brand['description']}")
                        st.markdown(f"**Analyses:** {len(get_analyses(brand['id']))}")
                    with dc2:
                        if st.button("✏️ Edit", key=f"be_{brand['id']}"): st.session_state[ekey]=True; st.rerun()
                        if st.button("🗑️ Delete", key=f"bd_{brand['id']}"): delete_brand(brand["id"]); st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
# ANALYZE
# ══════════════════════════════════════════════════════════════════════════════
elif page == "🔍 Analyze":
    st.markdown("<div class='hero'><h1>Analyze Strategy Deck</h1><p>Upload a deck to receive McKinsey-level strategy + visual design analysis</p></div>", unsafe_allow_html=True)

    brands = get_brands()
    if not brands: st.warning("Add at least one brand first."); st.stop()

    c1, c2 = st.columns([1,2])
    with c1:
        bmap = {b["name"]:b["id"] for b in brands}
        sname = st.selectbox("Select Brand", list(bmap.keys()))
        sbrand = get_brand(bmap[sname])
        st.markdown(f"""<div class="brand-row">
        <strong>{sbrand['name']}</strong><br>
        <span style="color:#64748B;font-size:.85rem">{sbrand['industry']}</span>
        {f"<br><span style='color:#94A3B8;font-size:.8rem'>{sbrand['target_market']}</span>" if sbrand.get('target_market') else ""}
        </div>""", unsafe_allow_html=True)
    with c2:
        uploaded = st.file_uploader("Upload Strategy Deck", type=["pdf","pptx"], help="PDF or PowerPoint. Max 50MB.")

    if uploaded:
        st.markdown(f"**📁 {uploaded.name}** · {uploaded.size/1024:.0f} KB")
        if st.button("🔍 Run Full Analysis", type="primary", use_container_width=True):
            with st.spinner("📖 Extracting content and visual metadata…"):
                try:
                    raw = uploaded.read()
                    text, visual_meta, image_samples = extract_text_and_visual(raw, uploaded.name)
                    if not text.strip(): st.error("No readable text found."); st.stop()
                    st.success(f"✅ Extracted {len(text.split())} words · {visual_meta.get('total_slides', visual_meta.get('total_pages','?'))} slides · {len(image_samples)} images found")
                except Exception as e: st.error(f"Parse error: {e}"); st.stop()

            with st.spinner("🧠 Running McKinsey-level strategy analysis…"):
                try:
                    strat_results = analyze_strategy(sbrand["name"], sbrand["industry"],
                        sbrand.get("description",""), sbrand.get("target_market",""), text, uploaded.name)
                except Exception as e: st.error(f"Strategy analysis error: {e}"); st.stop()

            with st.spinner("🎨 Analyzing visual design and deck quality…"):
                try:
                    visual_results = analyze_visual_appeal(visual_meta, image_samples,
                        sbrand["name"], sbrand["industry"], uploaded.name)
                except Exception as e: st.error(f"Visual analysis error: {e}"); st.stop()

            combined = round(strat_results["overall_score"] * 0.65 + visual_results["visual_score"] * 0.35)
            full_results = {"strategy": strat_results, "visual": visual_results, "combined_score": combined}

            save_analysis(sbrand["id"], uploaded.name, strat_results["overall_score"],
                visual_results["visual_score"], combined, strat_results["verdict"], full_results)
            st.success("✅ Analysis complete! Saved to History.")
            st.session_state["last_results"] = full_results

    if "last_results" in st.session_state:
        st.markdown("---")
        render_full_results(st.session_state["last_results"])


# ══════════════════════════════════════════════════════════════════════════════
# HISTORY
# ══════════════════════════════════════════════════════════════════════════════
elif page == "📋 History":
    st.markdown("<div class='hero'><h1>Analysis History</h1><p>Browse and compare all past strategy evaluations</p></div>", unsafe_allow_html=True)

    brands = get_brands()
    fopt = st.selectbox("Filter by Brand", ["All Brands"] + [b["name"] for b in brands])
    all_a = get_analyses() if fopt=="All Brands" else get_analyses(next(b["id"] for b in brands if b["name"]==fopt))

    if not all_a: st.info("No analyses found.")
    else:
        st.markdown(f"**{len(all_a)} result(s)**")
        for a in all_a:
            vs = a.get("visual_score") or "—"
            cs = a.get("combined_score") or a["overall_score"]
            with st.expander(f"📄  {a['deck_name']}  ·  {a['brand_name']}  ·  Combined: {cs}/100"):
                results = json.loads(a["results_json"])
                m1,m2,m3,m4 = st.columns(4)
                m1.metric("Combined Score", f"{cs}/100")
                m2.metric("Strategy Score", f"{a['overall_score']}/100")
                m3.metric("Visual Score", f"{vs}/100" if vs != '—' else "—")
                m4.markdown(f"**Verdict**<br><span style='color:{vc(a['verdict'])};font-weight:700'>{a['verdict']}</span>", unsafe_allow_html=True)

                if "strategy" in results and "visual" in results:
                    render_full_results(results)
                else:
                    st.markdown(results.get("strategy",{}).get("executive_summary","") or results.get("executive_summary",""))

                if st.button("🗑️ Delete", key=f"del_{a['id']}"): delete_analysis(a["id"]); st.rerun()
