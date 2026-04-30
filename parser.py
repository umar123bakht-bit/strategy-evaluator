import io
from pptx import Presentation
import pdfplumber


def extract_text_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def extract_text_from_pdf(file_bytes):
    parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                parts.append(f"[Page {i}]\n{text.strip()}")
    return "\n\n".join(parts)


def extract_text_from_file(file_bytes, filename):
    name = filename.lower()
    if name.endswith(".pptx"):
        return extract_text_from_pptx(file_bytes)
    elif name.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {filename}")
